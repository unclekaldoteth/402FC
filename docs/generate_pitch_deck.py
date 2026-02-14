from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_PATH = Path(__file__).resolve().parent / "402FC_Pitch_Deck.pptx"

# Brand and theme colors
COLOR_BG = RGBColor(9, 9, 11)
COLOR_SURFACE = RGBColor(24, 24, 27)
COLOR_TEXT = RGBColor(250, 250, 250)
COLOR_MUTED = RGBColor(161, 161, 170)
COLOR_ACCENT = RGBColor(252, 100, 50)  # Official Stacks orange extracted from logo asset
COLOR_SUCCESS = RGBColor(16, 185, 129)


def set_dark_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG


def add_top_bar(slide):
    bar = slide.shapes.add_shape(
        autoshape_type_id=1,  # Rectangle
        left=Inches(0.0),
        top=Inches(0.0),
        width=Inches(13.333),
        height=Inches(0.12),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_ACCENT
    bar.line.fill.background()


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.8), Inches(1.1))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = COLOR_TEXT

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.8))
        stf = sub_box.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        srun = sp.add_run()
        srun.text = subtitle
        srun.font.size = Pt(18)
        srun.font.color.rgb = COLOR_MUTED


def add_bullet_block(slide, bullets, left=0.9, top=2.2, width=11.5, height=4.6, font_size=24):
    box = slide.shapes.add_shape(
        autoshape_type_id=1,  # Rectangle
        left=Inches(left),
        top=Inches(top),
        width=Inches(width),
        height=Inches(height),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_SURFACE
    box.line.color.rgb = RGBColor(39, 39, 42)

    inner = slide.shapes.add_textbox(
        Inches(left + 0.45), Inches(top + 0.35), Inches(width - 0.9), Inches(height - 0.7)
    )
    tf = inner.text_frame
    tf.clear()
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.level = 0
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLOR_TEXT
        p.space_after = Pt(14)


def add_two_column_cards(slide, left_title, left_items, right_title, right_items):
    card_w = 5.75
    card_h = 4.4
    top = 2.2

    for idx, (title, items) in enumerate(
        [(left_title, left_items), (right_title, right_items)]
    ):
        left = 0.8 + idx * 6.0
        card = slide.shapes.add_shape(
            autoshape_type_id=1,
            left=Inches(left),
            top=Inches(top),
            width=Inches(card_w),
            height=Inches(card_h),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_SURFACE
        card.line.color.rgb = RGBColor(39, 39, 42)

        title_box = slide.shapes.add_textbox(Inches(left + 0.35), Inches(top + 0.25), Inches(5.2), Inches(0.6))
        ttf = title_box.text_frame
        ttf.clear()
        tp = ttf.paragraphs[0]
        tr = tp.add_run()
        tr.text = title
        tr.font.bold = True
        tr.font.size = Pt(22)
        tr.font.color.rgb = COLOR_ACCENT

        items_box = slide.shapes.add_textbox(
            Inches(left + 0.35), Inches(top + 0.95), Inches(5.2), Inches(3.2)
        )
        itf = items_box.text_frame
        itf.clear()
        itf.word_wrap = True

        for j, item in enumerate(items):
            p = itf.paragraphs[0] if j == 0 else itf.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = COLOR_TEXT
            p.space_after = Pt(10)


def add_metric_row(slide, metrics):
    y = 4.8
    card_w = 2.9
    gap = 0.35
    start_x = 0.9
    for i, (label, value) in enumerate(metrics):
        x = start_x + i * (card_w + gap)
        card = slide.shapes.add_shape(
            autoshape_type_id=1,
            left=Inches(x),
            top=Inches(y),
            width=Inches(card_w),
            height=Inches(1.65),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_SURFACE
        card.line.color.rgb = RGBColor(39, 39, 42)

        vbox = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.2), Inches(card_w - 0.4), Inches(0.65))
        vtf = vbox.text_frame
        vtf.clear()
        vp = vtf.paragraphs[0]
        vr = vp.add_run()
        vr.text = value
        vr.font.bold = True
        vr.font.size = Pt(26)
        vr.font.color.rgb = COLOR_ACCENT
        vp.alignment = PP_ALIGN.CENTER

        lbox = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.95), Inches(card_w - 0.4), Inches(0.5))
        ltf = lbox.text_frame
        ltf.clear()
        lp = ltf.paragraphs[0]
        lr = lp.add_run()
        lr.text = label
        lr.font.size = Pt(12)
        lr.font.color.rgb = COLOR_MUTED
        lp.alignment = PP_ALIGN.CENTER


def add_footer(slide, text):
    footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.8), Inches(0.35))
    tf = footer_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(12)
    run.font.color.rgb = COLOR_MUTED
    p.alignment = PP_ALIGN.RIGHT


def add_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_dark_background(slide)
    add_top_bar(slide)
    add_title(slide, title, subtitle)
    return slide


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide = add_slide(
        prs,
        "402FC: Pay-Per-Watch Football",
        "Making football access fairer for fans and more profitable for rights holders",
    )
    add_bullet_block(
        slide,
        [
            "Origin: Built from real fan pain in Indonesia",
            "Model: Pay only for the matches you actually watch",
            "Rail: x402 + STX micropayments for instant unlock",
            "Vision: Convert subscription-frustrated users into legal paying viewers",
        ],
        top=2.3,
        height=3.9,
        font_size=20,
    )
    add_footer(slide, "Hackathon pitch deck | 402FC")

    # Slide 2: Background
    slide = add_slide(prs, "Background", "Football passion is high, but access is rigid")
    add_bullet_block(
        slide,
        [
            "In Indonesia, football is one of the strongest fan cultures.",
            "Fans are loyal to one club and do not want to miss key matches.",
            "Many supporters only want selected fixtures, not full monthly bundles.",
            "When schedules conflict, users still pay the full subscription and feel wasted spend.",
        ],
    )
    add_footer(slide, "Context from fan behavior in Indonesia")

    # Slide 3: Problem
    slide = add_slide(prs, "Problem", "Subscription-only access creates friction and leakage")
    add_two_column_cards(
        slide,
        "User-side pain",
        [
            "Paying monthly for 1-2 matches feels inefficient",
            "Emergency schedules make paid subscriptions underused",
            "Price mismatch pushes fans to illegal streams",
        ],
        "Market-side loss",
        [
            "Rights holders miss casual and price-sensitive demand",
            "Piracy captures users willing to pay a fair per-match price",
            "Limited data on match-level willingness to pay",
        ],
    )
    add_footer(slide, "Core gap: pricing model does not match fan behavior")

    # Slide 4: Solution
    slide = add_slide(prs, "Solution", "402FC enables legal, flexible pay-per-watch access")
    add_bullet_block(
        slide,
        [
            "Per-event purchase: users pay only when they watch",
            "Micropayment unlock: wallet signs payment on demand",
            "No subscription lock-in, no long commitment",
            "Expandable paid modules: stream pass, highlights, analytics, AI insights",
        ],
    )
    add_metric_row(
        slide,
        [
            ("Entry Price", "0.02-0.08 STX"),
            ("Protocol", "HTTP 402"),
            ("Payment Rail", "x402 + STX"),
            ("Model", "Per Match"),
        ],
    )
    add_footer(slide, "Low-friction legal alternative to piracy")

    # Slide 5: Product flow
    slide = add_slide(prs, "How It Works", "Two-step payment flow")
    add_bullet_block(
        slide,
        [
            "1) User requests premium match content",
            "2) Server responds with HTTP 402 Payment Required",
            "3) Wallet signs micro STX payment",
            "4) Paid retry returns unlocked content/session",
        ],
        top=2.3,
        height=3.4,
        font_size=22,
    )
    add_bullet_block(
        slide,
        [
            "MVP already validates this preflight + paid retry behavior",
            "Integration tests lock flow reliability for future iterations",
        ],
        top=5.95,
        height=1.05,
        font_size=14,
    )
    add_footer(slide, "x402 payment challenge -> signed payment -> content unlock")

    # Slide 6: Market
    slide = add_slide(prs, "Market Opportunity", "Indonesia beachhead, then territory-by-territory expansion")
    add_two_column_cards(
        slide,
        "MVP market model",
        [
            "TAM: Large digital football audience",
            "SAM: Fans who prefer occasional event-based purchases",
            "SOM: Focused first-year capture via one-league launch",
        ],
        "Illustrative assumptions",
        [
            "Early paid capture: 0.2% to 0.8% of reachable SAM",
            "Paid events/user/month: 1 to 3",
            "Upside from high-intent derby and big-match windows",
        ],
    )
    add_footer(slide, "Illustrative planning model, not final financial guidance")

    # Slide 7: Benefits
    slide = add_slide(prs, "Stakeholder Benefits", "Win-win for fans and rights ecosystem")
    add_two_column_cards(
        slide,
        "For clubs/leagues/organizers",
        [
            "New incremental revenue from non-subscriber segments",
            "Monetize long-tail fixtures beyond marquee matches",
            "Reduce piracy pressure with fair legal pricing",
            "Get richer match-level demand and pricing data",
        ],
        "For fans",
        [
            "Pay only when they watch",
            "Lower commitment and clearer value per spend",
            "Faster, cleaner legal access for priority matches",
            "Optional premium modules without bundle lock-in",
        ],
    )
    add_footer(slide, "Business and consumer incentives are aligned")

    # Slide 8: Business model
    slide = add_slide(prs, "Monetization Model", "Flexible packaging around match-level intent")
    add_bullet_block(
        slide,
        [
            "Single match pass (core product)",
            "Derby pass and club-only mini bundles",
            "Replay/highlight unlocks after live window",
            "Premium analytics and AI prediction add-ons",
            "Future levers: sponsor-backed unlock credits and dynamic pricing",
        ],
    )
    add_metric_row(
        slide,
        [
            ("Primary KPI", "Paid Unlock Rate"),
            ("Retention KPI", "Repeat Purchases"),
            ("Ops KPI", "Playback Success"),
            ("Risk KPI", "Piracy Substitution"),
        ],
    )
    add_footer(slide, "Revenue scales with real match demand, not forced subscriptions")

    # Slide 9: Roadmap
    slide = add_slide(prs, "Execution Roadmap", "From hackathon MVP to licensed production")
    add_bullet_block(
        slide,
        [
            "Phase 1 (now): Stabilize payment + entitlement flow, polish UX",
            "Phase 2: Territory pilot with one licensed competition",
            "Phase 3: Add geo-rights, DRM, concurrency controls",
            "Phase 4: Expand by territory and add prediction intelligence",
        ],
        top=2.3,
        height=3.3,
        font_size=20,
    )
    milestone = slide.shapes.add_textbox(Inches(0.95), Inches(5.85), Inches(11.2), Inches(0.8))
    mtf = milestone.text_frame
    mtf.clear()
    p = mtf.paragraphs[0]
    run = p.add_run()
    run.text = "Near-term goal: prove conversion + reliability in one pilot market, then scale."
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = COLOR_SUCCESS
    add_footer(slide, "Rights-compliant expansion strategy")

    # Slide 10: Closing
    slide = add_slide(
        prs,
        "Closing Thesis",
        "402FC turns subscription-frustrated football fans into legal pay-per-watch customers",
    )
    add_bullet_block(
        slide,
        [
            "Fan truth: loyalty is club-centric, not always monthly-plan-centric.",
            "Product truth: match-level pricing creates a fair legal path.",
            "Business truth: rights holders gain incremental demand and better pricing intelligence.",
            "Ask: support pilot territory partnerships to validate this model at scale.",
        ],
        top=2.3,
        height=3.8,
        font_size=20,
    )
    cta = slide.shapes.add_textbox(Inches(0.9), Inches(6.25), Inches(11.5), Inches(0.7))
    ctf = cta.text_frame
    ctf.clear()
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run()
    cr.text = "Thank you | 402FC"
    cr.font.size = Pt(24)
    cr.font.bold = True
    cr.font.color.rgb = COLOR_ACCENT
    add_footer(slide, "Contact: 402FC project team")

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    print(f"Generated: {OUTPUT_PATH}")


if __name__ == "__main__":
    build_deck()
